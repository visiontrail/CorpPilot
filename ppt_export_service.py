"""
PowerPoint 导出服务
将 Dashboard 数据和图表导出为 PPT 演示文稿
"""

import base64
from io import BytesIO
from typing import Dict, List, Any
from pptx import Presentation
from pptx.chart.data import CategoryChartData, ChartData
from pptx.enum.chart import (
    XL_CHART_TYPE,
    XL_LEGEND_POSITION,
    XL_LABEL_POSITION
)
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from logger_config import get_logger


class PPTExporter:
    """PowerPoint 导出器"""

    def __init__(self):
        """初始化 PPT 导出器"""
        self.logger = get_logger("ppt_export_service")
        self.prs = Presentation()

        # 设置幻灯片尺寸为 16:9
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        self.logger.info("PPT 导出器初始化完成")

    def _decode_chart_image(self, base64_str: str) -> BytesIO:
        """
        将前端发送的 base64 图片解码为 BytesIO

        Args:
            base64_str: 格式 "data:image/png;base64,iVBORw0KG..."

        Returns:
            BytesIO 对象
        """
        try:
            # 移除 data:image/png;base64, 前缀
            if ',' in base64_str:
                image_data = base64_str.split(',')[1]
            else:
                image_data = base64_str

            # 解码 base64
            image_bytes = base64.b64decode(image_data)
            return BytesIO(image_bytes)
        except Exception as e:
            self.logger.error(f"图片解码失败: {str(e)}")
            raise

    def _add_title_slide(self, title: str, font_size: int = 32):
        """创建带标题的空白页，返回 slide 对象"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(font_size)
        title_para.font.bold = True
        title_para.font.name = "微软雅黑"
        title_para.font.color.rgb = RGBColor(38, 38, 38)
        title_para.alignment = PP_ALIGN.CENTER

        return slide

    def create_kpi_slide(self, kpi_data: Dict[str, Any]):
        """
        创建 KPI 指标幻灯片

        Args:
            kpi_data: KPI 数据字典
        """
        self.logger.debug("开始创建 KPI 指标幻灯片")

        # 使用空白布局
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = "核心指标概览"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.name = "微软雅黑"
        title_para.font.color.rgb = RGBColor(38, 38, 38)
        title_para.alignment = PP_ALIGN.CENTER

        # KPI 数据
        kpis = [
            {
                "label": "总成本",
                "value": f"¥{kpi_data.get('total_cost', 0):,.2f}",
                "color": RGBColor(82, 196, 26)  # 绿色
            },
            {
                "label": "订单总数",
                "value": str(kpi_data.get('total_orders', 0)),
                "color": RGBColor(24, 144, 255)  # 蓝色
            },
            {
                "label": "异常记录",
                "value": str(kpi_data.get('anomaly_count', 0)),
                "color": RGBColor(245, 34, 45)  # 红色
            },
            {
                "label": "超标订单",
                "value": str(kpi_data.get('over_standard_count', 0)),
                "color": RGBColor(250, 140, 22)  # 橙色
            },
            {
                "label": "紧急预订率",
                "value": f"{kpi_data.get('urgent_booking_ratio', 0):.1f}%",
                "color": RGBColor(114, 46, 209)  # 紫色
            }
        ]

        # 布局：3个在上排，2个在下排
        positions = [
            (1.5, 2.0),   # 第1个
            (5.0, 2.0),   # 第2个
            (8.5, 2.0),   # 第3个
            (2.75, 4.5),  # 第4个
            (6.75, 4.5),  # 第5个
        ]

        box_width = 3.0
        box_height = 1.8

        for i, (kpi, pos) in enumerate(zip(kpis, positions)):
            # 创建 KPI 卡片背景
            shape = slide.shapes.add_shape(
                1,  # 矩形
                Inches(pos[0]), Inches(pos[1]),
                Inches(box_width), Inches(box_height)
            )

            # 设置背景颜色（浅色）
            fill = shape.fill
            fill.solid()
            color = kpi['color']
            r, g, b = int(color[0]), int(color[1]), int(color[2])
            fill.fore_color.rgb = RGBColor(
                min(255, r + 200),
                min(255, g + 200),
                min(255, b + 200)
            )

            # 设置边框
            line = shape.line
            line.color.rgb = kpi['color']
            line.width = Pt(2)

            # 添加标签文本
            label_box = slide.shapes.add_textbox(
                Inches(pos[0]), Inches(pos[1] + 0.3),
                Inches(box_width), Inches(0.5)
            )
            label_frame = label_box.text_frame
            label_frame.text = kpi['label']
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(20)
            label_para.font.name = "微软雅黑"
            label_para.font.color.rgb = RGBColor(89, 89, 89)
            label_para.alignment = PP_ALIGN.CENTER

            # 添加数值文本
            value_box = slide.shapes.add_textbox(
                Inches(pos[0]), Inches(pos[1] + 0.9),
                Inches(box_width), Inches(0.7)
            )
            value_frame = value_box.text_frame
            value_frame.text = kpi['value']
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(32)
            value_para.font.bold = True
            value_para.font.name = "微软雅黑"
            value_para.font.color.rgb = kpi['color']
            value_para.alignment = PP_ALIGN.CENTER

        self.logger.info("KPI 指标幻灯片创建完成")

    def create_chart_slide(self, title: str, chart_image_base64: str):
        """
        创建图表幻灯片

        Args:
            title: 幻灯片标题
            chart_image_base64: 图表的 base64 图片数据
        """
        self.logger.debug(f"开始创建图表幻灯片: {title}")

        # 使用空白布局
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.name = "微软雅黑"
        title_para.font.color.rgb = RGBColor(38, 38, 38)
        title_para.alignment = PP_ALIGN.CENTER

        # 解码并添加图表图片
        try:
            image_stream = self._decode_chart_image(chart_image_base64)

            # 添加图片到幻灯片
            # 位置: 左边距0.5英寸, 顶部1.5英寸（标题下方）
            # 宽度: 12.333英寸（几乎全宽）, 高度: 5.5英寸
            slide.shapes.add_picture(
                image_stream,
                Inches(0.5), Inches(1.5),
                width=Inches(12.333)
            )

            self.logger.info(f"图表幻灯片创建完成: {title}")
        except Exception as e:
            self.logger.error(f"添加图表图片失败: {str(e)}")
            raise

    def create_department_cost_chart(self, dept_metrics: List[Dict[str, Any]]):
        """
        使用部门成本数据创建内置饼图（环形图）

        Args:
            dept_metrics: 部门指标列表，包含"一级部门"、"总成本"等字段
        """
        if not dept_metrics:
            self.logger.warning("部门成本数据为空，跳过部门成本图表页")
            return

        self.logger.debug("开始创建部门成本环形图幻灯片")
        slide = self._add_title_slide("部门成本分布", font_size=32)

        # 仅取前15个部门以保持可读性
        categories = []
        values = []
        for dept in dept_metrics[:15]:
            categories.append(str(dept.get('一级部门', '未知') or '未知'))
            values.append(float(dept.get('总成本', 0) or 0))

        if not any(values):
            self.logger.warning("部门成本数据全部为0，跳过图表绘制")
            return

        chart_data = ChartData()
        chart_data.categories = categories
        chart_data.add_series('部门成本', values)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(0.5), Inches(1.2),
            Inches(12.333), Inches(5.8),
            chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT

        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_percentage = True
        data_labels.show_category_name = True
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(10)
        data_labels.font.name = "微软雅黑"

        self.logger.info("部门成本环形图幻灯片创建完成")

    def create_project_cost_chart(self, projects: List[Dict[str, Any]]):
        """
        使用项目成本数据创建横向柱状图

        Args:
            projects: 项目成本列表，包含"项目代码"、"项目名称"、"总成本"等字段
        """
        if not projects:
            self.logger.warning("项目成本数据为空，跳过项目成本柱状图")
            return

        self.logger.debug("开始创建项目成本柱状图幻灯片")
        slide = self._add_title_slide("项目成本排名（Top 20）", font_size=32)

        categories = []
        values = []
        for project in projects[:20]:
            code = str(project.get('项目代码', '') or '未知')
            name = str(project.get('项目名称', '') or '')
            label = f"{code} {name}".strip()[:40] if name else code
            categories.append(label)
            values.append(float(project.get('总成本', 0) or 0))

        if not any(values):
            self.logger.warning("项目成本数据全部为0，跳过图表绘制")
            return

        chart_data = CategoryChartData()
        chart_data.categories = categories
        chart_data.add_series('总成本', values)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            Inches(0.5), Inches(1.2),
            Inches(12.333), Inches(5.8),
            chart_data
        ).chart

        chart.has_legend = False
        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.number_format = '#,##0'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(9)
        data_labels.font.name = "微软雅黑"

        category_axis = chart.category_axis
        category_axis.tick_labels.font.size = Pt(9)
        category_axis.tick_labels.font.name = "微软雅黑"

        value_axis = chart.value_axis
        value_axis.tick_labels.number_format = '#,##0'

        self.logger.info("项目成本柱状图幻灯片创建完成")

    def create_over_standard_chart(self, over_stats: Dict[str, Any], total_orders: int):
        """
        创建超标占比环形图

        Args:
            over_stats: 超标统计字典，包含 total/flight/hotel/train
            total_orders: 总订单数（用于计算合规订单）
        """
        if over_stats is None:
            self.logger.warning("超标统计数据为空，跳过超标占比图表")
            return

        over_total = int(over_stats.get('total', 0) or 0)
        total_orders = int(total_orders or 0)
        compliant = max(total_orders - over_total, 0)

        # 若总订单未知但有超标数据，至少展示超标占比
        if total_orders == 0 and over_total > 0:
            compliant = 0
            total_orders = over_total

        if over_total == 0 and compliant == 0:
            self.logger.warning("超标与合规订单均为0，跳过超标占比图表")
            return

        slide = self._add_title_slide("超标订单占比", font_size=32)

        chart_data = ChartData()
        chart_data.categories = ['超标订单', '合规订单']
        chart_data.add_series('订单数', [over_total, compliant])

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(0.5), Inches(1.2),
            Inches(12.333), Inches(5.8),
            chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT

        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_percentage = True
        data_labels.show_category_name = True
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(10)
        data_labels.font.name = "微软雅黑"

        self.logger.info("超标占比环形图幻灯片创建完成")

    def create_booking_behavior_chart(self, booking_behavior: Dict[str, Any]):
        """
        创建预订行为占比环形图（紧急预订 vs 正常预订）

        Args:
            booking_behavior: 预订行为统计字典，包含 total_orders/urgent_orders
        """
        if not booking_behavior:
            self.logger.warning("预订行为数据为空，跳过预订行为图表")
            return

        total_orders = int(booking_behavior.get('total_orders', 0) or 0)
        urgent_orders = int(booking_behavior.get('urgent_orders', 0) or 0)
        normal_orders = max(total_orders - urgent_orders, 0)

        if urgent_orders == 0 and normal_orders == 0:
            self.logger.warning("预订行为订单为0，跳过预订行为图表")
            return

        slide = self._add_title_slide("预订行为概览", font_size=32)

        chart_data = ChartData()
        chart_data.categories = ['紧急预订(≤2天)', '正常预订']
        chart_data.add_series('订单数', [urgent_orders, normal_orders])

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(0.5), Inches(1.2),
            Inches(12.333), Inches(5.8),
            chart_data
        ).chart

        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT

        plot = chart.plots[0]
        plot.has_data_labels = True
        data_labels = plot.data_labels
        data_labels.show_percentage = True
        data_labels.show_category_name = True
        data_labels.number_format = '0%'
        data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        data_labels.font.size = Pt(10)
        data_labels.font.name = "微软雅黑"

        self.logger.info("预订行为环形图幻灯片创建完成")

    def create_table_slide(self, title: str, headers: List[str], rows: List[List[Any]]):
        """
        创建数据表格幻灯片

        Args:
            title: 幻灯片标题
            headers: 表头列表
            rows: 数据行列表
        """
        self.logger.debug(f"开始创建表格幻灯片: {title}, 行数: {len(rows)}")

        # 使用空白布局
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.name = "微软雅黑"
        title_para.font.color.rgb = RGBColor(38, 38, 38)
        title_para.alignment = PP_ALIGN.CENTER

        # 计算表格尺寸
        num_cols = len(headers)
        num_rows = min(len(rows), 20) + 1  # 限制最多20行数据 + 1行表头

        # 添加表格
        table = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(0.5), Inches(1.2),
            Inches(12.333), Inches(5.8)
        ).table

        # 设置表头
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)

            # 表头样式
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.bold = True
                paragraph.font.name = "微软雅黑"
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.alignment = PP_ALIGN.CENTER

            # 表头背景色（蓝色）
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(68, 114, 196)

        # 填充数据行
        for row_idx, row_data in enumerate(rows[:20]):  # 限制最多20行
            for col_idx, cell_value in enumerate(row_data):
                if col_idx >= num_cols:
                    break

                cell = table.cell(row_idx + 1, col_idx)

                # 处理数值格式化
                if isinstance(cell_value, float):
                    cell.text = f"{cell_value:,.2f}"
                elif isinstance(cell_value, int):
                    cell.text = f"{cell_value:,}"
                else:
                    # 截断过长文本
                    text = str(cell_value)
                    cell.text = text[:100] + "..." if len(text) > 100 else text

                # 数据单元格样式
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.font.name = "微软雅黑"
                    paragraph.alignment = PP_ALIGN.CENTER

                # 交替行背景色
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)

        self.logger.info(f"表格幻灯片创建完成: {title}, 实际行数: {num_rows}")

    def export_to_bytes(self) -> BytesIO:
        """
        将 PPT 导出为字节流

        Returns:
            BytesIO 对象
        """
        self.logger.debug("开始导出 PPT 到字节流")

        output_stream = BytesIO()
        self.prs.save(output_stream)
        output_stream.seek(0)

        self.logger.info("PPT 导出到字节流完成")
        return output_stream
